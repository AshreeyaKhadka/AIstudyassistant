from flask import Blueprint, request, jsonify
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from config import db
from models.content import StudentUpload
from services.auth_service import login_required
import os
import time
import logging

upload_bp = Blueprint('upload', __name__)
logger = logging.getLogger(__name__)

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# 15 MB Max file size
MAX_FILE_SIZE = 15 * 1024 * 1024
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx'}

def get_file_extension(filename):
    if '.' not in filename:
        return ''
    return filename.rsplit('.', 1)[1].lower()

def extract_text_from_file(filepath, ext):
    text = ""
    if ext == 'pdf':
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text() or ""
        doc.close()
    elif ext == 'docx':
        doc = docx.Document(filepath)
        text_list = []
        for p in doc.paragraphs:
            if p.text:
                text_list.append(p.text)
        text = "\n".join(text_list)
    elif ext == 'pptx':
        prs = Presentation(filepath)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_list.append(shape.text)
        text = "\n".join(text_list)
    return text

@upload_bp.route('/', methods=['POST'])
@login_required
def upload_file(user):
    # Check if a file part is present
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
        
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    ext = get_file_extension(file.filename)
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": "Invalid file type. Only PDF, DOCX, and PPTX are supported."}), 400

    # Retrieve and validate subject
    subject = request.form.get('subject', '').strip()
    if not subject:
        subject = "General"

    # Enforce maximum upload count limit (15 files)
    upload_count = StudentUpload.query.filter_by(user_id=user.id).count()
    if upload_count >= 15:
        return jsonify({"error": "Upload limit of 15 materials reached. Please delete some before uploading new ones."}), 403

    # Generate secure, unique filename to avoid duplicates
    orig_filename = file.filename
    sec_name = secure_filename(orig_filename)
    filename = f"{user.id}_{int(time.time())}_{sec_name}"
    filepath = os.path.join(UPLOAD_FOLDER, filename)

    try:
        # Save file to disk
        file.save(filepath)
        size_bytes = os.path.getsize(filepath)
    except Exception as e:
        logger.error(f"Failed to save uploaded file: {e}")
        return jsonify({"error": "Failed to save file to server storage"}), 500

    # Validate file size limits after saving
    if size_bytes > MAX_FILE_SIZE:
        try:
            os.remove(filepath)
        except Exception:
            pass
        return jsonify({"error": "File size exceeds the 15 MB limit."}), 400

    # Parse and extract text content based on file type
    parsed_text = ""
    try:
        parsed_text = extract_text_from_file(filepath, ext)
    except Exception as e:
        logger.error(f"Failed to extract text from {ext} file: {e}")
        try:
            os.remove(filepath)
        except Exception:
            pass
        return jsonify({"error": f"Failed to parse document text: {str(e)}"}), 500

    # Save to the database
    upload_record = StudentUpload(
        user_id=user.id,
        filename=filename,
        original_filename=orig_filename,
        file_type=ext,
        file_size=size_bytes,
        subject=subject,
        storage_path=filepath,
        parsed_text=parsed_text
    )

    try:
        db.session.add(upload_record)
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        logger.error(f"Database save error: {e}")
        try:
            os.remove(filepath)
        except Exception:
            pass
        return jsonify({"error": "Failed to persist document record in database"}), 500

    return jsonify({
        "message": "Material uploaded and parsed successfully",
        "upload": upload_record.to_dict(),
        "parsed_preview": parsed_text[:200] if parsed_text else ""
    }), 200

@upload_bp.route('/', methods=['GET'])
@login_required
def get_uploads(user):
    try:
        # Query materials uploaded by the current student
        uploads = StudentUpload.query.filter_by(user_id=user.id).order_by(StudentUpload.created_at.desc()).all()
        return jsonify([u.to_dict() for u in uploads]), 200
    except Exception as e:
        logger.error(f"Database error fetching uploads: {e}")
        return jsonify({"error": "Failed to fetch uploaded materials"}), 500

@upload_bp.route('/<int:upload_id>', methods=['DELETE'])
@login_required
def delete_upload(user, upload_id):
    try:
        # Fetch the target upload
        upload = StudentUpload.query.filter_by(id=upload_id, user_id=user.id).first()
        if not upload:
            return jsonify({"error": "Material not found or access unauthorized"}), 404

        # Delete local file from disk storage
        if os.path.exists(upload.storage_path):
            try:
                os.remove(upload.storage_path)
            except Exception as e:
                logger.error(f"Failed to delete storage file {upload.storage_path}: {e}")

        # Delete database record
        db.session.delete(upload)
        db.session.commit()

        return jsonify({"message": "Material deleted successfully", "id": upload_id}), 200
    except Exception as e:
        db.session.rollback()
        logger.error(f"Error during material deletion: {e}")
        return jsonify({"error": "Failed to delete study material"}), 500
