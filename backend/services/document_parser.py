import base64
import logging
import mimetypes
import os
import re

import fitz
import requests

from config import Config

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {'.pdf', '.txt', '.png', '.jpg', '.jpeg', '.webp', '.pptx'}
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp'}
OCR_MIN_ALNUM_CHARS = 20
OCR_MIN_IMAGE_COVERAGE = 0.25
MAX_OCR_PAGES = 25


def is_supported_material(filename: str) -> bool:
    _, ext = os.path.splitext(filename or '')
    return ext.lower() in SUPPORTED_EXTENSIONS


def supported_material_message() -> str:
    return 'Supported files: PDF slides/notes, PPTX slides, TXT notes, and PNG/JPG/WEBP handwritten-note images.'


def parse_uploaded_material(file, filepath: str) -> str:
    text, _metadata = parse_uploaded_material_with_metadata(file, filepath)
    return text


def parse_uploaded_material_with_metadata(file, filepath: str, enable_ocr: bool = True) -> tuple[str, dict]:
    """Compatibility wrapper that saves an upload before extracting it."""
    filename = file.filename or ''
    if not is_supported_material(filename):
        raise ValueError(supported_material_message())
    file.save(filepath)
    return extract_material_from_path(filepath, filename, enable_ocr=enable_ocr)


def extract_material_from_path(filepath: str, filename: str = '', enable_ocr: bool = True) -> tuple[str, dict]:
    """Extract structured study text and diagnostics from a file already on disk."""
    _, ext = os.path.splitext(filename or filepath)
    ext = ext.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(supported_material_message())

    if ext == '.txt':
        with open(filepath, 'r', encoding='utf-8', errors='replace') as source:
            raw_text = source.read().strip()
        warnings = []
        if '\ufffd' in raw_text:
            warnings.append('Some text characters could not be decoded exactly.')
        text = f'[Page 1]\n{raw_text}' if raw_text else ''
        return text, _extraction_metadata('typed_text', text, 1, warnings)

    if ext == '.pdf':
        return _parse_pdf(filepath, enable_ocr=enable_ocr)

    if ext == '.pptx':
        return _parse_pptx(filepath)

    if ext in IMAGE_EXTENSIONS:
        text = _ocr_image_file(filepath).strip()
        marked_text = f'[Page 1]\n{text}' if text else ''
        return marked_text, _extraction_metadata('ocr', marked_text, 1, [])

    raise ValueError(supported_material_message())


def _extraction_metadata(
    method: str, text: str, page_count: int, warnings: list[str], empty_pages: int = 0,
    native_text_pages: int = 0, ocr_pages: int = 0, image_pages: int = 0,
) -> dict:
    character_count = len((text or '').strip())
    content_characters = max(0, character_count - max(1, page_count) * 10)
    density = content_characters / max(1, page_count)
    empty_ratio = empty_pages / max(1, page_count)

    if character_count == 0 or density < 50 or empty_ratio > 0.5:
        quality = 'low'
    elif density < 180 or empty_ratio > 0.2 or warnings:
        quality = 'partial'
    else:
        quality = 'good'

    return {
        'extraction_method': method,
        'extraction_quality': quality,
        'character_count': character_count,
        'page_count': page_count,
        'empty_pages': empty_pages,
        'native_text_pages': native_text_pages,
        'ocr_pages': ocr_pages,
        'image_pages': image_pages,
        'warnings': list(dict.fromkeys(warnings)),
    }


def _page_image_coverage(page) -> float:
    page_area = max(1.0, float(page.rect.width * page.rect.height))
    covered_area = 0.0
    try:
        for image in page.get_image_info(xrefs=False):
            bbox = image.get('bbox')
            if not bbox or len(bbox) != 4:
                continue
            width = max(0.0, float(bbox[2]) - float(bbox[0]))
            height = max(0.0, float(bbox[3]) - float(bbox[1]))
            covered_area += width * height
    except Exception:
        return 0.0
    return min(1.0, covered_area / page_area)


def _should_ocr_pdf_page(page, extracted_text: str) -> bool:
    """OCR only image-backed pages whose selectable text is effectively unusable."""
    alnum_count = len(re.findall(r'[A-Za-z0-9]', extracted_text or ''))
    if alnum_count >= OCR_MIN_ALNUM_CHARS:
        return False
    return _page_image_coverage(page) >= OCR_MIN_IMAGE_COVERAGE


def _open_pdf(filepath: str):
    try:
        document = fitz.open(filepath)
        if document.needs_pass:
            document.close()
            raise ValueError('This PDF is password-protected. Please upload an unprotected version.')
        return document
    except ValueError:
        raise
    except Exception as exc:
        message = str(exc).lower()
        if 'password' in message or 'encrypted' in message:
            raise ValueError('This PDF is password-protected. Please upload an unprotected version.') from exc
        if any(term in message for term in ('cannot open', 'broken', 'truncat')):
            raise ValueError('This PDF appears to be corrupted or incomplete. Please try a different file.') from exc
        raise ValueError(f'Could not open this PDF: {exc}') from exc


def _parse_pdf(filepath: str, enable_ocr: bool = True) -> tuple[str, dict]:
    document = _open_pdf(filepath)
    sections = []
    warnings = []
    empty_pages = 0
    ocr_attempts = 0
    ocr_successes = 0
    native_text_pages = 0
    image_pages = 0
    forced_ocr = False

    try:
        page_count = len(document)
        for index, page in enumerate(document):
            page_number = index + 1
            extracted = (page.get_text('text', sort=True) or '').strip()
            selected_text = extracted
            if _page_image_coverage(page) >= OCR_MIN_IMAGE_COVERAGE:
                image_pages += 1

            if enable_ocr and _should_ocr_pdf_page(page, extracted):
                if ocr_attempts >= MAX_OCR_PAGES:
                    warnings.append(f'OCR limit reached; page {page_number} was not OCR processed.')
                elif not Config.GEMINI_API_KEY:
                    warnings.append(f'Page {page_number} has little selectable text and OCR is not configured.')
                else:
                    ocr_attempts += 1
                    try:
                        pixmap = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
                        ocr_text = _ocr_image_bytes(pixmap.tobytes('png'), 'image/png').strip()
                        if len(ocr_text) > len(selected_text):
                            selected_text = ocr_text
                            ocr_successes += 1
                    except Exception as exc:
                        logger.warning('OCR failed for PDF page %s: %s', page_number, exc)
                        warnings.append(f'OCR failed for page {page_number}; selectable text was kept where available.')

            if selected_text:
                sections.append(f'[Page {page_number}]\n{selected_text}')
                if selected_text == extracted:
                    native_text_pages += 1
            else:
                empty_pages += 1
                warnings.append(f'No readable text was found on page {page_number}.')

        # Some scanned PDFs contain full-page vector or mask content that image
        # coverage detection cannot identify. Only force OCR when the complete
        # native/selective pass produced no usable text.
        if enable_ocr and not sections and page_count and Config.GEMINI_API_KEY:
            forced_ocr = True
            warnings.append('No selectable text was found; full-page OCR was attempted.')
            for index, page in enumerate(document):
                if ocr_attempts >= MAX_OCR_PAGES:
                    warnings.append(f'OCR limit reached; pages after {MAX_OCR_PAGES} were not processed.')
                    break
                ocr_attempts += 1
                try:
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                    ocr_text = _ocr_image_bytes(pixmap.tobytes('png'), 'image/png').strip()
                    if ocr_text:
                        sections.append(f'[Page {index + 1}]\n{ocr_text}')
                        ocr_successes += 1
                except Exception as exc:
                    logger.warning('Forced OCR failed for PDF page %s: %s', index + 1, exc)
                    warnings.append(f'OCR failed for page {index + 1}.')
        elif enable_ocr and not sections and page_count and not Config.GEMINI_API_KEY:
            warnings.append('This PDF appears scanned, but OCR is not configured on the server.')
        elif not enable_ocr and not sections and page_count:
            if image_pages:
                warnings.append('This PDF contains scanned page images and no selectable text.')
            else:
                warnings.append('No selectable text was found in this PDF.')
    finally:
        document.close()

    method = 'pdf_forced_ocr' if forced_ocr and ocr_successes else 'pdf_text_ocr' if ocr_successes else 'pdf_text'
    text = '\n\n'.join(sections)
    return text, _extraction_metadata(
        method, text, page_count, warnings, empty_pages,
        native_text_pages=native_text_pages,
        ocr_pages=ocr_successes,
        image_pages=image_pages,
    )


def _parse_pptx(filepath: str) -> tuple[str, dict]:
    try:
        from pptx import Presentation
        presentation = Presentation(filepath)
    except ImportError as exc:
        raise RuntimeError('PPTX parsing requires python-pptx. Install backend requirements first.') from exc
    except Exception as exc:
        raise ValueError('This slide deck appears to be corrupted or is not a valid PPTX file.') from exc

    slides = []
    empty_slides = 0
    warnings = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        seen = set()
        for shape in slide.shapes:
            shape_text = (getattr(shape, 'text', '') or '').strip()
            if shape_text and shape_text not in seen:
                text_parts.append(shape_text)
                seen.add(shape_text)
            if getattr(shape, 'has_table', False):
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text and row_text not in seen:
                        text_parts.append(row_text)
                        seen.add(row_text)
        if text_parts:
            slides.append(f'[Slide {index}]\n' + '\n'.join(text_parts))
        else:
            empty_slides += 1

    if empty_slides:
        warnings.append(f'{empty_slides} slide(s) contained no extractable text.')
    text = '\n\n'.join(slides)
    return text, _extraction_metadata('slide_text', text, len(presentation.slides), warnings, empty_slides)


def _ocr_image_file(filepath: str) -> str:
    mime_type = mimetypes.guess_type(filepath)[0] or 'image/png'
    with open(filepath, 'rb') as source:
        return _ocr_image_bytes(source.read(), mime_type)


def _ocr_image_bytes(image_bytes: bytes, mime_type: str) -> str:
    if not Config.GEMINI_API_KEY:
        raise RuntimeError('GEMINI_API_KEY is required for OCR of handwritten or scanned notes.')

    payload = {
        'contents': [{
            'role': 'user',
            'parts': [
                {
                    'text': (
                        'Extract all readable study-note text from this image. Preserve headings, '
                        'bullet points, formulas, table rows, code, and slide structure. Return plain '
                        'text only. Mark unclear fragments as [unclear] instead of inventing text.'
                    )
                },
                {
                    'inline_data': {
                        'mime_type': mime_type,
                        'data': base64.b64encode(image_bytes).decode('ascii'),
                    }
                },
            ],
        }],
        'generationConfig': {'temperature': 0.0, 'maxOutputTokens': 4096},
    }

    try:
        response = requests.post(
            f"{Config.GEMINI_API_BASE_URL.rstrip('/')}/models/{Config.GEMINI_MODEL}:generateContent",
            headers={'x-goog-api-key': Config.GEMINI_API_KEY},
            json=payload,
            timeout=90,
        )
        response.raise_for_status()
    except requests.RequestException as exc:
        logger.error('Gemini OCR request failed: %s', exc)
        raise RuntimeError(f'OCR failed: {exc}') from exc

    try:
        data = response.json()
    except ValueError as exc:
        raise RuntimeError('OCR service returned an invalid response.') from exc
    candidates = data.get('candidates') or []
    if not candidates:
        return ''
    content = (candidates[0] or {}).get('content') or {}
    parts = content.get('parts') or []
    return '\n'.join(part.get('text', '') for part in parts if isinstance(part, dict)).strip()
